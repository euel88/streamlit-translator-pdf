#!/usr/bin/env python3
"""
문서 형식별 처리 모듈
각 문서 형식에 대한 번역 처리를 담당합니다.
Streamlit Cloud 호환 버전
"""

import os
import sys
import shutil
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from abc import ABC, abstractmethod
import json
import re
from datetime import datetime
import tempfile
import fitz  # PyMuPDF

logger = logging.getLogger(__name__)


class DocumentHandler(ABC):
    """문서 처리기 추상 클래스"""
    
    def __init__(self, translator):
        self.translator = translator
        
    @abstractmethod
    def translate(self, file_path: Path, source_lang: str, target_lang: str,
                 service: str, output_dir: Path, preserve_formatting: bool,
                 enable_ocr: bool, callback: Optional[callable] = None) -> Dict:
        """문서 번역 실행"""
        pass
    
    def _create_output_filename(self, input_path: Path, suffix: str = "") -> str:
        """출력 파일명 생성"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        base_name = input_path.stem
        extension = input_path.suffix
        
        if suffix:
            return f"{base_name}_{suffix}_{timestamp}{extension}"
        else:
            return f"{base_name}_translated_{timestamp}{extension}"
    
    def _should_translate(self, text: str) -> bool:
        """번역이 필요한 텍스트인지 확인"""
        if not text or not text.strip():
            return False
        
        # 숫자나 특수문자만 있는 경우 제외
        if re.match(r'^[\d\s\W]+$', text):
            return False
        
        # 너무 짧은 텍스트 제외 (단, CJK 문자는 예외)
        if len(text.strip()) < 3 and not re.search(r'[\u4e00-\u9fff\u3040-\u309f\u30a0-\u30ff\uac00-\ud7af]', text):
            return False
        
        return True
    
    def _safe_get_ocr_engine(self):
        """OCR 엔진을 안전하게 가져오기"""
        try:
            from ocr_engine import get_ocr_engine
            ocr = get_ocr_engine()
            if ocr.is_available():
                return ocr
            else:
                logger.info("OCR 엔진을 사용할 수 없습니다.")
                return None
        except ImportError:
            logger.info("OCR 모듈이 설치되지 않았습니다.")
            return None
        except Exception as e:
            logger.warning(f"OCR 엔진 로드 실패: {e}")
            return None


class PDFHandler(DocumentHandler):
    """PDF 문서 처리기 - PyMuPDF 직접 사용"""
    
    def translate(self, file_path: Path, source_lang: str, target_lang: str,
                 service: str, output_dir: Path, preserve_formatting: bool,
                 enable_ocr: bool, callback: Optional[callable] = None) -> Dict:
        """PDF 번역 실행 - PyMuPDF 사용"""
        logger.info(f"PDF 번역 시작: {file_path.name}")
        
        try:
            # PyMuPDF로 PDF 열기
            pdf_document = fitz.open(str(file_path))
            
            # 새 PDF 생성
            output_pdf = fitz.open()
            
            total_pages = len(pdf_document)
            
            for page_num in range(total_pages):
                if callback:
                    callback(page_num + 1, total_pages, f"페이지 {page_num + 1}/{total_pages} 번역 중...")
                
                page = pdf_document[page_num]
                
                # 텍스트 추출
                text_dict = page.get_text("dict")
                
                # 페이지 크기 가져오기
                page_rect = page.rect
                
                # 새 페이지 생성
                new_page = output_pdf.new_page(width=page_rect.width, height=page_rect.height)
                
                # 배경 이미지 복사 (있는 경우)
                pix = page.get_pixmap(alpha=False)
                img_data = pix.tobytes("png")
                new_page.insert_image(page_rect, stream=img_data)
                
                # 텍스트 블록 번역 및 삽입
                for block in text_dict["blocks"]:
                    if block["type"] == 0:  # 텍스트 블록
                        for line in block["lines"]:
                            for span in line["spans"]:
                                original_text = span["text"]
                                
                                if self._should_translate(original_text):
                                    try:
                                        # 텍스트 번역
                                        translated_text = self.translator.translate_text(
                                            original_text, source_lang, target_lang, service
                                        )
                                        
                                        # 번역된 텍스트 삽입
                                        font_size = span["size"]
                                        flags = span["flags"]
                                        
                                        # 기존 텍스트 영역에 흰색 사각형 그리기 (텍스트 지우기)
                                        bbox = span["bbox"]
                                        rect = fitz.Rect(bbox)
                                        new_page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1))
                                        
                                        # 번역된 텍스트 삽입
                                        point = fitz.Point(bbox[0], bbox[3] - 2)  # 좌하단 기준
                                        
                                        # 폰트 설정
                                        fontname = "helv"  # 기본 폰트
                                        if "font" in span:
                                            if "cjk" in span["font"].lower() or "korean" in span["font"].lower():
                                                fontname = "china-s"  # CJK 폰트
                                        
                                        try:
                                            new_page.insert_text(
                                                point,
                                                translated_text,
                                                fontsize=font_size,
                                                fontname=fontname,
                                                color=(0, 0, 0)
                                            )
                                        except:
                                            # 폰트 오류 시 기본 폰트 사용
                                            new_page.insert_text(
                                                point,
                                                translated_text,
                                                fontsize=font_size,
                                                fontname="helv",
                                                color=(0, 0, 0)
                                            )
                                    
                                    except Exception as e:
                                        logger.warning(f"텍스트 번역 실패: {e}")
                                        # 원본 텍스트 유지
                                        point = fitz.Point(span["bbox"][0], span["bbox"][3] - 2)
                                        new_page.insert_text(
                                            point,
                                            original_text,
                                            fontsize=span["size"],
                                            color=(0, 0, 0)
                                        )
            
            # PDF 저장
            output_path = output_dir / self._create_output_filename(file_path)
            output_pdf.save(str(output_path))
            
            # 리소스 정리
            output_pdf.close()
            pdf_document.close()
            
            # 대조본 생성 (선택사항)
            dual_path = None
            if preserve_formatting:
                try:
                    dual_path = self._create_dual_pdf(file_path, output_path, output_dir)
                except Exception as e:
                    logger.warning(f"대조본 생성 실패: {e}")
            
            return {
                'output_file': str(output_path),
                'dual_file': str(dual_path) if dual_path else None,
                'format': 'pdf',
                'pages': total_pages
            }
            
        except Exception as e:
            logger.error(f"PDF 번역 실패: {e}")
            
            # 폴백: 간단한 텍스트 추출 및 번역
            try:
                return self._fallback_pdf_translation(
                    file_path, source_lang, target_lang, service, output_dir
                )
            except Exception as fallback_error:
                logger.error(f"폴백 번역도 실패: {fallback_error}")
                raise
    
    def _fallback_pdf_translation(self, file_path: Path, source_lang: str, 
                                 target_lang: str, service: str, output_dir: Path) -> Dict:
        """폴백: 단순 텍스트 추출 후 새 PDF 생성"""
        logger.info("폴백 PDF 번역 모드 사용")
        
        try:
            # PDF에서 텍스트 추출
            pdf_document = fitz.open(str(file_path))
            
            # 텍스트 수집
            all_text = []
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                text = page.get_text()
                if text.strip():
                    all_text.append(f"[Page {page_num + 1}]\n{text}")
            
            full_text = "\n\n".join(all_text)
            
            # 전체 텍스트 번역
            translated_text = ""
            if full_text:
                # 텍스트를 청크로 분할 (너무 길면 API 제한에 걸릴 수 있음)
                chunks = self._split_text(full_text, 3000)
                translated_chunks = []
                
                for chunk in chunks:
                    if self._should_translate(chunk):
                        try:
                            translated_chunk = self.translator.translate_text(
                                chunk, source_lang, target_lang, service
                            )
                            translated_chunks.append(translated_chunk)
                        except Exception as e:
                            logger.warning(f"청크 번역 실패: {e}")
                            translated_chunks.append(chunk)
                    else:
                        translated_chunks.append(chunk)
                
                translated_text = "\n\n".join(translated_chunks)
            
            # 새 PDF 생성 (텍스트만)
            output_pdf = fitz.open()
            
            # A4 크기 페이지 생성
            page_width = 595  # A4 width in points
            page_height = 842  # A4 height in points
            
            # 텍스트를 페이지에 맞게 분할
            lines = translated_text.split('\n')
            current_page = output_pdf.new_page(width=page_width, height=page_height)
            y_position = 50
            line_height = 14
            margin = 50
            
            for line in lines:
                if y_position > page_height - margin:
                    # 새 페이지 생성
                    current_page = output_pdf.new_page(width=page_width, height=page_height)
                    y_position = 50
                
                # 텍스트 삽입
                if line.strip():
                    try:
                        current_page.insert_text(
                            fitz.Point(margin, y_position),
                            line[:100],  # 긴 줄은 자르기
                            fontsize=11,
                            fontname="helv",
                            color=(0, 0, 0)
                        )
                    except:
                        pass
                
                y_position += line_height
            
            # PDF 저장
            output_path = output_dir / self._create_output_filename(file_path, "text")
            output_pdf.save(str(output_path))
            
            # 리소스 정리
            output_pdf.close()
            pdf_document.close()
            
            return {
                'output_file': str(output_path),
                'dual_file': None,
                'format': 'pdf',
                'pages': len(output_pdf),
                'note': '텍스트 전용 번역 (레이아웃 미보존)'
            }
            
        except Exception as e:
            logger.error(f"폴백 PDF 번역 실패: {e}")
            raise
    
    def _split_text(self, text: str, max_length: int) -> List[str]:
        """텍스트를 청크로 분할"""
        chunks = []
        current_chunk = ""
        
        sentences = text.split('. ')
        for sentence in sentences:
            if len(current_chunk) + len(sentence) < max_length:
                current_chunk += sentence + '. '
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + '. '
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _create_dual_pdf(self, original_path: Path, translated_path: Path, 
                        output_dir: Path) -> Path:
        """원문-번역 대조본 PDF 생성"""
        try:
            original_pdf = fitz.open(str(original_path))
            translated_pdf = fitz.open(str(translated_path))
            dual_pdf = fitz.open()
            
            # 각 페이지를 나란히 배치
            for i in range(min(len(original_pdf), len(translated_pdf))):
                orig_page = original_pdf[i]
                trans_page = translated_pdf[i]
                
                # 두 페이지를 합친 크기의 새 페이지 생성
                width = orig_page.rect.width + trans_page.rect.width
                height = max(orig_page.rect.height, trans_page.rect.height)
                
                new_page = dual_pdf.new_page(width=width, height=height)
                
                # 원본 페이지 삽입 (왼쪽)
                new_page.show_pdf_page(
                    fitz.Rect(0, 0, orig_page.rect.width, height),
                    original_pdf,
                    i
                )
                
                # 번역 페이지 삽입 (오른쪽)
                new_page.show_pdf_page(
                    fitz.Rect(orig_page.rect.width, 0, width, height),
                    translated_pdf,
                    i
                )
            
            # 저장
            dual_path = output_dir / self._create_output_filename(original_path, "dual")
            dual_pdf.save(str(dual_path))
            
            # 리소스 정리
            dual_pdf.close()
            translated_pdf.close()
            original_pdf.close()
            
            return dual_path
            
        except Exception as e:
            logger.error(f"대조본 생성 실패: {e}")
            return None


class WordHandler(DocumentHandler):
    """Word 문서 처리기"""
    
    def translate(self, file_path: Path, source_lang: str, target_lang: str,
                 service: str, output_dir: Path, preserve_formatting: bool,
                 enable_ocr: bool, callback: Optional[callable] = None) -> Dict:
        """Word 문서 번역"""
        logger.info(f"Word 번역 시작: {file_path.name}")
        
        try:
            from docx import Document
            from docx.shared import RGBColor, Pt
            from docx.enum.text import WD_COLOR_INDEX
            
            # 문서 열기
            doc = Document(str(file_path))
            output_doc = Document()
            
            # 문서 속성 복사
            if preserve_formatting:
                output_doc.core_properties.author = doc.core_properties.author
                output_doc.core_properties.title = doc.core_properties.title
            
            total_items = len(doc.paragraphs) + sum(len(table.rows) * len(table.columns) for table in doc.tables)
            processed_items = 0
            
            # 단락 번역
            for paragraph in doc.paragraphs:
                if callback:
                    processed_items += 1
                    callback(processed_items, total_items, "단락 번역 중...")
                
                # 빈 단락 처리
                if not paragraph.text.strip():
                    output_doc.add_paragraph()
                    continue
                
                # 새 단락 생성
                new_para = output_doc.add_paragraph()
                
                # 서식 복사
                if preserve_formatting:
                    new_para.alignment = paragraph.alignment
                    new_para.paragraph_format.space_before = paragraph.paragraph_format.space_before
                    new_para.paragraph_format.space_after = paragraph.paragraph_format.space_after
                    new_para.paragraph_format.line_spacing = paragraph.paragraph_format.line_spacing
                
                # Run 단위로 번역 (서식 보존)
                for run in paragraph.runs:
                    if self._should_translate(run.text):
                        try:
                            translated_text = self.translator.translate_text(
                                run.text, source_lang, target_lang, service
                            )
                        except Exception as e:
                            logger.warning(f"텍스트 번역 실패: {e}")
                            translated_text = run.text
                    else:
                        translated_text = run.text
                    
                    # 새 Run 생성
                    new_run = new_para.add_run(translated_text)
                    
                    # 서식 복사
                    if preserve_formatting:
                        if run.bold is not None:
                            new_run.bold = run.bold
                        if run.italic is not None:
                            new_run.italic = run.italic
                        if run.underline:
                            new_run.underline = run.underline
                        if run.font.size:
                            new_run.font.size = run.font.size
                        if run.font.name:
                            new_run.font.name = run.font.name
                        if run.font.color.rgb:
                            new_run.font.color.rgb = run.font.color.rgb
            
            # 표 번역
            for table_idx, table in enumerate(doc.tables):
                if callback:
                    callback(processed_items, total_items, f"표 {table_idx + 1} 번역 중...")
                
                # 새 표 생성
                new_table = output_doc.add_table(rows=len(table.rows), cols=len(table.columns))
                
                # 표 스타일 복사
                if preserve_formatting and table.style:
                    new_table.style = table.style
                
                # 셀 번역
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        processed_items += 1
                        
                        new_cell = new_table.rows[row_idx].cells[col_idx]
                        
                        # 셀 텍스트 번역
                        cell_text = cell.text.strip()
                        if self._should_translate(cell_text):
                            try:
                                translated_text = self.translator.translate_text(
                                    cell_text, source_lang, target_lang, service
                                )
                                new_cell.text = translated_text
                            except Exception as e:
                                logger.warning(f"셀 번역 실패: {e}")
                                new_cell.text = cell_text
                        else:
                            new_cell.text = cell_text
            
            # 헤더/푸터 번역
            for section in doc.sections:
                # 헤더
                if section.header:
                    for paragraph in section.header.paragraphs:
                        if self._should_translate(paragraph.text):
                            try:
                                translated = self.translator.translate_text(
                                    paragraph.text, source_lang, target_lang, service
                                )
                                # 새 섹션에 헤더 추가
                                if output_doc.sections:
                                    output_doc.sections[-1].header.paragraphs[0].text = translated
                            except:
                                pass
                
                # 푸터
                if section.footer:
                    for paragraph in section.footer.paragraphs:
                        if self._should_translate(paragraph.text):
                            try:
                                translated = self.translator.translate_text(
                                    paragraph.text, source_lang, target_lang, service
                                )
                                if output_doc.sections:
                                    output_doc.sections[-1].footer.paragraphs[0].text = translated
                            except:
                                pass
            
            # 파일 저장
            output_path = output_dir / self._create_output_filename(file_path)
            output_doc.save(str(output_path))
            
            # 대조본 생성 (원본 + 번역)
            dual_path = output_dir / self._create_output_filename(file_path, "dual")
            self._create_word_dual(doc, output_doc, dual_path)
            
            return {
                'output_file': str(output_path),
                'dual_file': str(dual_path),
                'format': 'docx',
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables)
            }
            
        except Exception as e:
            logger.error(f"Word 번역 실패: {e}")
            raise
    
    def _create_word_dual(self, original_doc, translated_doc, output_path: Path):
        """원문-번역 대조본 생성"""
        try:
            from docx import Document
            
            dual_doc = Document()
            dual_doc.add_heading('원문-번역 대조본', 0)
            dual_doc.add_paragraph(f'생성일: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
            dual_doc.add_page_break()
            
            # 단락별 대조
            for orig_para, trans_para in zip(original_doc.paragraphs, translated_doc.paragraphs):
                if orig_para.text.strip():
                    # 원문
                    p1 = dual_doc.add_paragraph()
                    p1.add_run('[원문] ').bold = True
                    p1.add_run(orig_para.text)
                    
                    # 번역
                    p2 = dual_doc.add_paragraph()
                    p2.add_run('[번역] ').bold = True
                    p2.add_run(trans_para.text)
                    
                    dual_doc.add_paragraph()  # 빈 줄
            
            dual_doc.save(str(output_path))
            
        except Exception as e:
            logger.warning(f"대조본 생성 실패: {e}")


class ExcelHandler(DocumentHandler):
    """Excel 문서 처리기"""
    
    def translate(self, file_path: Path, source_lang: str, target_lang: str,
                 service: str, output_dir: Path, preserve_formatting: bool,
                 enable_ocr: bool, callback: Optional[callable] = None) -> Dict:
        """Excel 문서 번역"""
        logger.info(f"Excel 번역 시작: {file_path.name}")
        
        try:
            from openpyxl import load_workbook, Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
            
            # 워크북 열기
            wb = load_workbook(str(file_path), data_only=False)
            output_wb = Workbook()
            
            # 기본 시트 제거
            if 'Sheet' in output_wb.sheetnames:
                output_wb.remove(output_wb['Sheet'])
            
            total_cells = sum(ws.max_row * ws.max_column for ws in wb.worksheets)
            processed_cells = 0
            
            # 시트별 처리
            for sheet_idx, ws in enumerate(wb.worksheets):
                if callback:
                    callback(processed_cells, total_cells, f"시트 '{ws.title}' 번역 중...")
                
                # 새 시트 생성
                output_ws = output_wb.create_sheet(title=ws.title)
                
                # 열 너비 복사
                if preserve_formatting:
                    for col in ws.column_dimensions:
                        output_ws.column_dimensions[col].width = ws.column_dimensions[col].width
                
                # 행 높이 복사
                if preserve_formatting:
                    for row in ws.row_dimensions:
                        output_ws.row_dimensions[row].height = ws.row_dimensions[row].height
                
                # 셀 처리
                for row in range(1, ws.max_row + 1):
                    for col in range(1, ws.max_column + 1):
                        processed_cells += 1
                        
                        cell = ws.cell(row=row, column=col)
                        output_cell = output_ws.cell(row=row, column=col)
                        
                        # 수식 처리
                        if cell.data_type == 'f':
                            # 수식은 그대로 복사
                            output_cell.value = cell.value
                            output_cell._value = cell._value
                            output_cell.data_type = cell.data_type
                        
                        # 텍스트 처리
                        elif cell.value and isinstance(cell.value, str):
                            if self._should_translate(cell.value):
                                try:
                                    translated_text = self.translator.translate_text(
                                        cell.value, source_lang, target_lang, service
                                    )
                                    output_cell.value = translated_text
                                except Exception as e:
                                    logger.warning(f"셀 번역 실패 ({row},{col}): {e}")
                                    output_cell.value = cell.value
                            else:
                                output_cell.value = cell.value
                        else:
                            # 숫자, 날짜 등은 그대로 복사
                            output_cell.value = cell.value
                        
                        # 서식 복사
                        if preserve_formatting:
                            # 폰트
                            if cell.font:
                                output_cell.font = Font(
                                    name=cell.font.name,
                                    size=cell.font.size,
                                    bold=cell.font.bold,
                                    italic=cell.font.italic,
                                    color=cell.font.color
                                )
                            
                            # 배경색
                            if cell.fill:
                                output_cell.fill = PatternFill(
                                    fill_type=cell.fill.fill_type,
                                    start_color=cell.fill.start_color,
                                    end_color=cell.fill.end_color
                                )
                            
                            # 정렬
                            if cell.alignment:
                                output_cell.alignment = Alignment(
                                    horizontal=cell.alignment.horizontal,
                                    vertical=cell.alignment.vertical,
                                    wrap_text=cell.alignment.wrap_text
                                )
                            
                            # 테두리
                            if cell.border:
                                output_cell.border = cell.border
                            
                            # 숫자 서식
                            if cell.number_format:
                                output_cell.number_format = cell.number_format
                
                # 병합 셀 처리
                for merged_range in ws.merged_cells.ranges:
                    output_ws.merge_cells(str(merged_range))
                
                # 차트 제목 번역
                for chart in ws._charts:
                    if chart.title and self._should_translate(str(chart.title)):
                        try:
                            translated_title = self.translator.translate_text(
                                str(chart.title), source_lang, target_lang, service
                            )
                            # 차트는 복사가 복잡하므로 제목만 번역
                            logger.info(f"차트 제목 번역: {chart.title} → {translated_title}")
                        except:
                            pass
            
            # 파일 저장
            output_path = output_dir / self._create_output_filename(file_path)
            output_wb.save(str(output_path))
            
            # 주석 추가본 생성
            annotated_path = output_dir / self._create_output_filename(file_path, "annotated")
            self._create_excel_annotated(wb, output_wb, annotated_path)
            
            return {
                'output_file': str(output_path),
                'annotated_file': str(annotated_path),
                'format': 'xlsx',
                'sheets': len(wb.worksheets),
                'total_cells': processed_cells
            }
            
        except Exception as e:
            logger.error(f"Excel 번역 실패: {e}")
            raise
    
    def _create_excel_annotated(self, original_wb, translated_wb, output_path: Path):
        """주석이 추가된 버전 생성 (원문을 주석으로)"""
        try:
            from openpyxl.comments import Comment
            from openpyxl import load_workbook
            
            # 번역본 복사
            translated_wb.save(str(output_path))
            annotated_wb = load_workbook(str(output_path))
            
            # 각 시트별로 주석 추가
            for orig_ws, anno_ws in zip(original_wb.worksheets, annotated_wb.worksheets):
                for row in range(1, min(orig_ws.max_row + 1, 100)):  # 최대 100행
                    for col in range(1, min(orig_ws.max_column + 1, 20)):  # 최대 20열
                        orig_cell = orig_ws.cell(row=row, column=col)
                        anno_cell = anno_ws.cell(row=row, column=col)
                        
                        # 원문이 있고 번역된 경우 주석 추가
                        if (orig_cell.value and 
                            isinstance(orig_cell.value, str) and 
                            anno_cell.value != orig_cell.value):
                            comment = Comment(f"원문: {orig_cell.value}", "Translator")
                            anno_cell.comment = comment
            
            annotated_wb.save(str(output_path))
            
        except Exception as e:
            logger.warning(f"주석 추가본 생성 실패: {e}")


class PowerPointHandler(DocumentHandler):
    """PowerPoint 문서 처리기"""
    
    def translate(self, file_path: Path, source_lang: str, target_lang: str,
                 service: str, output_dir: Path, preserve_formatting: bool,
                 enable_ocr: bool, callback: Optional[callable] = None) -> Dict:
        """PowerPoint 프레젠테이션 번역"""
        logger.info(f"PowerPoint 번역 시작: {file_path.name}")
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # 프레젠테이션 열기
            prs = Presentation(str(file_path))
            output_prs = Presentation()
            
            # 슬라이드 크기 복사
            output_prs.slide_width = prs.slide_width
            output_prs.slide_height = prs.slide_height
            
            total_items = sum(
                len(slide.shapes) + len(getattr(slide, 'notes_slide', []))
                for slide in prs.slides
            )
            processed_items = 0
            
            # 슬라이드별 처리
            for slide_idx, slide in enumerate(prs.slides):
                if callback:
                    callback(processed_items, total_items, f"슬라이드 {slide_idx + 1} 번역 중...")
                
                # 슬라이드 레이아웃 복사
                slide_layout = output_prs.slide_layouts[min(slide_idx, len(output_prs.slide_layouts) - 1)]
                output_slide = output_prs.slides.add_slide(slide_layout)
                
                # 도형별 처리
                for shape in slide.shapes:
                    processed_items += 1
                    
                    # 텍스트 프레임이 있는 경우
                    if shape.has_text_frame:
                        # 해당하는 도형 찾기 또는 생성
                        output_shape = self._find_or_create_shape(output_slide, shape)
                        
                        # 단락별 번역
                        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                            if self._should_translate(paragraph.text):
                                try:
                                    translated_text = self.translator.translate_text(
                                        paragraph.text, source_lang, target_lang, service
                                    )
                                    
                                    # 출력 단락 설정
                                    if para_idx < len(output_shape.text_frame.paragraphs):
                                        output_para = output_shape.text_frame.paragraphs[para_idx]
                                    else:
                                        output_para = output_shape.text_frame.add_paragraph()
                                    
                                    output_para.text = translated_text
                                    
                                    # 서식 복사
                                    if preserve_formatting and paragraph.runs:
                                        run = paragraph.runs[0]
                                        output_run = output_para.runs[0] if output_para.runs else output_para.add_run()
                                        
                                        if run.font.name:
                                            output_run.font.name = run.font.name
                                        if run.font.size:
                                            output_run.font.size = run.font.size
                                        if run.font.bold is not None:
                                            output_run.font.bold = run.font.bold
                                        if run.font.italic is not None:
                                            output_run.font.italic = run.font.italic
                                    
                                except Exception as e:
                                    logger.warning(f"텍스트 번역 실패: {e}")
                                    if para_idx < len(output_shape.text_frame.paragraphs):
                                        output_shape.text_frame.paragraphs[para_idx].text = paragraph.text
                    
                    # 표가 있는 경우
                    elif shape.has_table:
                        table = shape.table
                        # 표 번역 로직
                        for row in table.rows:
                            for cell in row.cells:
                                if self._should_translate(cell.text):
                                    try:
                                        translated = self.translator.translate_text(
                                            cell.text, source_lang, target_lang, service
                                        )
                                        cell.text = translated
                                    except:
                                        pass
                    
                    # 차트가 있는 경우
                    elif shape.has_chart:
                        chart = shape.chart
                        # 차트 제목 번역
                        if chart.has_title and self._should_translate(chart.chart_title.text_frame.text):
                            try:
                                translated = self.translator.translate_text(
                                    chart.chart_title.text_frame.text,
                                    source_lang, target_lang, service
                                )
                                # 차트는 복사가 복잡하므로 로그만
                                logger.info(f"차트 제목: {chart.chart_title.text_frame.text} → {translated}")
                            except:
                                pass
                
                # 발표자 노트 번역
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if self._should_translate(notes_text):
                        try:
                            translated_notes = self.translator.translate_text(
                                notes_text, source_lang, target_lang, service
                            )
                            # 노트 추가
                            if output_slide.has_notes_slide:
                                output_slide.notes_slide.notes_text_frame.text = translated_notes
                        except Exception as e:
                            logger.warning(f"노트 번역 실패: {e}")
            
            # 파일 저장
            output_path = output_dir / self._create_output_filename(file_path)
            output_prs.save(str(output_path))
            
            return {
                'output_file': str(output_path),
                'format': 'pptx',
                'slides': len(prs.slides),
                'shapes': processed_items
            }
            
        except Exception as e:
            logger.error(f"PowerPoint 번역 실패: {e}")
            raise
    
    def _find_or_create_shape(self, slide, original_shape):
        """슬라이드에서 해당하는 도형 찾기 또는 생성"""
        # 간단한 구현 - 실제로는 더 정교한 매칭 필요
        for shape in slide.shapes:
            if shape.shape_type == original_shape.shape_type:
                return shape
        
        # 새 텍스트 상자 생성
        from pptx.util import Inches
        left = top = Inches(1)
        width = height = Inches(3)
        
        return slide.shapes.add_textbox(left, top, width, height)


class TextHandler(DocumentHandler):
    """텍스트 파일 처리기 (txt, md)"""
    
    def translate(self, file_path: Path, source_lang: str, target_lang: str,
                 service: str, output_dir: Path, preserve_formatting: bool,
                 enable_ocr: bool, callback: Optional[callable] = None) -> Dict:
        """텍스트 파일 번역"""
        logger.info(f"텍스트 파일 번역 시작: {file_path.name}")
        
        try:
            # 파일 읽기
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Markdown 파일인 경우
            if file_path.suffix.lower() == '.md':
                translated_content = self._translate_markdown(
                    content, source_lang, target_lang, service, preserve_formatting
                )
            else:
                # 일반 텍스트 파일
                lines = content.split('\n')
                translated_lines = []
                
                total_lines = len(lines)
                for idx, line in enumerate(lines):
                    if callback:
                        callback(idx + 1, total_lines, f"줄 {idx + 1}/{total_lines} 번역 중...")
                    
                    if self._should_translate(line):
                        try:
                            translated_line = self.translator.translate_text(
                                line, source_lang, target_lang, service
                            )
                            translated_lines.append(translated_line)
                        except Exception as e:
                            logger.warning(f"줄 번역 실패: {e}")
                            translated_lines.append(line)
                    else:
                        translated_lines.append(line)
                
                translated_content = '\n'.join(translated_lines)
            
            # 파일 저장
            output_path = output_dir / self._create_output_filename(file_path)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(translated_content)
            
            # 대조본 생성
            dual_path = output_dir / self._create_output_filename(file_path, "dual")
            self._create_text_dual(content, translated_content, dual_path)
            
            return {
                'output_file': str(output_path),
                'dual_file': str(dual_path),
                'format': file_path.suffix[1:],
                'lines': len(content.split('\n')),
                'characters': len(content)
            }
            
        except Exception as e:
            logger.error(f"텍스트 파일 번역 실패: {e}")
            raise
    
    def _translate_markdown(self, content: str, source_lang: str, target_lang: str,
                          service: str, preserve_formatting: bool) -> str:
        """Markdown 문서 번역 (구조 보존)"""
        import re
        
        # 코드 블록 보호
        code_blocks = []
        code_pattern = r'```[\s\S]*?```'
        
        def save_code_block(match):
            code_blocks.append(match.group(0))
            return f"__CODE_BLOCK_{len(code_blocks)-1}__"
        
        content = re.sub(code_pattern, save_code_block, content)
        
        # 인라인 코드 보호
        inline_codes = []
        inline_pattern = r'`[^`]+`'
        
        def save_inline_code(match):
            inline_codes.append(match.group(0))
            return f"__INLINE_CODE_{len(inline_codes)-1}__"
        
        content = re.sub(inline_pattern, save_inline_code, content)
        
        # 링크 URL 보호
        links = []
        link_pattern = r'\[([^\]]+)\]\(([^)]+)\)'
        
        def process_link(match):
            text = match.group(1)
            url = match.group(2)
            links.append(url)
            
            # 링크 텍스트만 번역
            if self._should_translate(text):
                try:
                    translated_text = self.translator.translate_text(
                        text, source_lang, target_lang, service
                    )
                    return f"[{translated_text}](__LINK_{len(links)-1}__)"
                except:
                    return f"[{text}](__LINK_{len(links)-1}__)"
            else:
                return f"[{text}](__LINK_{len(links)-1}__)"
        
        content = re.sub(link_pattern, process_link, content)
        
        # 줄별 번역
        lines = content.split('\n')
        translated_lines = []
        
        for line in lines:
            # 빈 줄이나 특수 마크다운 요소는 그대로
            if not line.strip() or line.strip().startswith(('---', '***', '___')):
                translated_lines.append(line)
                continue
            
            # 제목 처리
            heading_match = re.match(r'^(#{1,6})\s+(.+)', line)
            if heading_match:
                level = heading_match.group(1)
                text = heading_match.group(2)
                
                if self._should_translate(text):
                    try:
                        translated_text = self.translator.translate_text(
                            text, source_lang, target_lang, service
                        )
                        translated_lines.append(f"{level} {translated_text}")
                    except:
                        translated_lines.append(line)
                else:
                    translated_lines.append(line)
                continue
            
            # 리스트 항목 처리
            list_match = re.match(r'^(\s*[-*+]|\s*\d+\.)\s+(.+)', line)
            if list_match:
                prefix = list_match.group(1)
                text = list_match.group(2)
                
                if self._should_translate(text):
                    try:
                        translated_text = self.translator.translate_text(
                            text, source_lang, target_lang, service
                        )
                        translated_lines.append(f"{prefix} {translated_text}")
                    except:
                        translated_lines.append(line)
                else:
                    translated_lines.append(line)
                continue
            
            # 일반 텍스트
            if self._should_translate(line):
                try:
                    translated_line = self.translator.translate_text(
                        line, source_lang, target_lang, service
                    )
                    translated_lines.append(translated_line)
                except:
                    translated_lines.append(line)
            else:
                translated_lines.append(line)
        
        # 번역된 내용 합치기
        translated_content = '\n'.join(translated_lines)
        
        # 보호된 요소 복원
        # 링크 복원
        for idx, url in enumerate(links):
            translated_content = translated_content.replace(f"__LINK_{idx}__", url)
        
        # 인라인 코드 복원
        for idx, code in enumerate(inline_codes):
            translated_content = translated_content.replace(f"__INLINE_CODE_{idx}__", code)
        
        # 코드 블록 복원
        for idx, block in enumerate(code_blocks):
            translated_content = translated_content.replace(f"__CODE_BLOCK_{idx}__", block)
        
        return translated_content
    
    def _create_text_dual(self, original: str, translated: str, output_path: Path):
        """텍스트 대조본 생성"""
        try:
            orig_lines = original.split('\n')
            trans_lines = translated.split('\n')
            
            dual_content = []
            dual_content.append("=" * 50)
            dual_content.append("원문-번역 대조본")
            dual_content.append(f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            dual_content.append("=" * 50)
            dual_content.append("")
            
            # 줄별 대조
            max_lines = max(len(orig_lines), len(trans_lines))
            for i in range(max_lines):
                if i < len(orig_lines) and orig_lines[i].strip():
                    dual_content.append(f"[원문] {orig_lines[i]}")
                    
                if i < len(trans_lines) and trans_lines[i].strip():
                    dual_content.append(f"[번역] {trans_lines[i]}")
                
                dual_content.append("")  # 빈 줄
            
            # 파일 저장
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(dual_content))
                
        except Exception as e:
            logger.warning(f"대조본 생성 실패: {e}")
